import openai
import pptx

client = openai()

# Set your OpenAI API key


def extract_text_from_pptx(file_path):
    presentation = pptx.Presentation(file_path)
    text = ""
    for slide_number, slide in enumerate(presentation.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

def chat_with_gpt(prompt):
    response = client.chat.Completions.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": prompt},
        ]
    )
    return response['choices'][0]['message']['content'].strip()

def main():
    # Replace with the path to your PPTX file
    pptx_file_path = "C:/Users/Khalil Ur Rehman/Desktop/PPtx AI Sumerizer/idea about campaign final.pptx"

    # Extract text from PPTX file
    pptx_text = extract_text_from_pptx(pptx_file_path)

    # Provide a summary of the PPTX content
    summary_prompt = f"Summarize the content of the PPTX file:\n{pptx_text}"
    summary_response = chat_with_gpt(summary_prompt)
    print("Summary:", summary_response)

    # Allow the user to ask questions about the PPTX content
    while True:
        user_question = input("Ask a question (type 'exit' to end): ")
        if user_question.lower() == 'exit':
            break
        question_prompt = f"Q: {user_question}\nA:"
        answer_response = chat_with_gpt(question_prompt)
        print("Answer:", answer_response)

if __name__ == "__main__":
    main()
